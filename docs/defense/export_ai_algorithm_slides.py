from __future__ import annotations

import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PPTX_DIR = BASE_DIR / "pptx"
PDF_DIR = BASE_DIR / "pdf"
NOTES_DIR = BASE_DIR / "notes"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG = RGBColor(246, 243, 237)
PANEL = RGBColor(255, 252, 247)
INK = RGBColor(30, 33, 31)
MUTED = RGBColor(88, 95, 91)
GREEN = RGBColor(41, 91, 67)
GREEN_LIGHT = RGBColor(215, 232, 220)
AMBER = RGBColor(198, 129, 55)
AMBER_LIGHT = RGBColor(244, 229, 211)
LINE = RGBColor(205, 198, 186)


@dataclass(frozen=True)
class DeckCopy:
    lang: str
    deck_title: str
    deck_subtitle: str
    slide_titles: List[str]
    title_scope: List[str]
    problem_cards: List[tuple[str, str]]
    pipeline_labels: List[str]
    pipeline_caption: str
    filter_chips: List[str]
    filter_note: str
    feedback_bullets: List[str]
    feedback_note: str
    current_now: List[str]
    next_phase: List[str]
    closing_statement: str
    notes_title: str
    speaker_notes: List[str]


COPY_BY_LANG = {
    "en": DeckCopy(
        lang="en",
        deck_title="OpenNutri AI / Algorithm Layer",
        deck_subtitle="Midterm exposition | Retrieval, ranking, and feedback reuse",
        slide_titles=[
            "AI / Algorithm Layer",
            "Why This Layer Matters",
            "Retrieval And Ranking Pipeline",
            "Implemented Feedback Loop",
            "Current State And Scope Boundary",
        ],
        title_scope=[
            "Multi-source paper discovery",
            "Bilingual EN/TR retrieval strategy",
            "Feedback-driven re-ranking for later runs",
        ],
        problem_cards=[
            (
                "Distributed literature",
                "Relevant food-composition papers are spread across multiple databases and local platforms.",
            ),
            (
                "High acquisition cost",
                "Downloading every possible PDF is slow, noisy, and wastes expert review time.",
            ),
            (
                "Retrieval design",
                "The system ranks metadata before acquisition and treats English and Turkish as separate pools.",
            ),
        ],
        pipeline_labels=[
            "Search\nSources",
            "Search\nGate",
            "Metadata\nFilter",
            "PDF\nIntake",
            "Expert\nAnnotation",
        ],
        pipeline_caption="The pipeline rejects weak candidates early and reserves PDF download for stronger papers.",
        filter_chips=[
            "composition phrases",
            "food and nutrient terms",
            "unit patterns",
            "semantic similarity",
            "source and query priors",
            "feedback-weighted terms",
        ],
        filter_note="Current implementation uses additive scoring: explicit lexical cues, embeddings, and learned feedback all contribute.",
        feedback_bullets=[
            "Latest visible label per user",
            "Positive, negative, and conflict sets",
            "Updated query and anchor phrases",
            "Updated concept and source priorities",
        ],
        feedback_note="Current state: a working batch feedback loop. It changes later runs, but it is not yet a trained classifier.",
        current_now=[
            "Multi-source search across Europe PMC, OpenAlex, Semantic Scholar, and DergiPark",
            "Two-stage filtering before PDF download",
            "Embedding-supported metadata scoring",
            "Feedback-driven statistical term updates",
            "Paper-stock refill that closes the operational loop",
        ],
        next_phase=[
            "Document segmentation",
            "LLM-assisted extraction",
            "Classifier training after label volume grows",
        ],
        closing_statement="Key midterm claim: retrieval, annotation, and feedback already work together as one system.",
        notes_title="OpenNutri AI / Algorithm Mini Deck Notes",
        speaker_notes=[
            "Slide 1: Frame your role as retrieval, ranking, bilingual search, and feedback reuse. Keep the introduction under 30 seconds.",
            "Slide 2: Stress that the problem is not scraping alone. The real issue is deciding which papers deserve download and expert time.",
            "Slide 3: Explain the staged pipeline slowly. The strongest phrase here is 'Search -> Filter -> Acquisition'.",
            "Slide 4: Say clearly that the feedback loop is implemented, but batch-updated rather than online. This is one of the most important distinctions.",
            "Slide 5: Be explicit about boundaries. Say what works now, then state what is deferred to the second semester without sounding apologetic.",
        ],
    ),
    "tr": DeckCopy(
        lang="tr",
        deck_title="OpenNutri Yapay Zeka / Algoritma Katmanı",
        deck_subtitle="Arasınav sunumu | Erişim, sıralama ve geri besleme kullanımı",
        slide_titles=[
            "Yapay Zeka / Algoritma Katmanı",
            "Bu Katman Neden Gerekli",
            "Erişim Ve Sıralama Hattı",
            "Uygulanmış Geri Besleme Döngüsü",
            "Mevcut Durum Ve Kapsam Sınırı",
        ],
        title_scope=[
            "Çok kaynaklı makale keşfi",
            "EN/TR ayrımlı erişim stratejisi",
            "Sonraki çalıştırmalar için geri besleme tabanlı yeniden sıralama",
        ],
        problem_cards=[
            (
                "Dağınık literatür",
                "Gıda bileşimiyle ilgili uygun makaleler birden fazla veritabanı ve yerel platforma dağılmış durumdadır.",
            ),
            (
                "Yüksek edinim maliyeti",
                "Olası her PDF'yi indirmek yavaştır, gürültülüdür ve uzman inceleme zamanını boşa harcar.",
            ),
            (
                "Erişim tasarımı",
                "Sistem PDF ediniminden önce meta veriyi sıralar ve İngilizce ile Türkçeyi ayrı havuzlar olarak ele alır.",
            ),
        ],
        pipeline_labels=[
            "Arama\nKaynakları",
            "Arama\nKapısı",
            "Meta Veri\nFiltresi",
            "PDF\nEdinimi",
            "Uzman\nAnotasyonu",
        ],
        pipeline_caption="Hat, zayıf adayları erken eler ve PDF indirmeyi daha güçlü adaylara ayırır.",
        filter_chips=[
            "bileşim ifadeleri",
            "gıda ve besin terimleri",
            "birim örüntüleri",
            "anlamsal benzerlik",
            "kaynak ve sorgu öncelikleri",
            "geri besleme ağırlıklı terimler",
        ],
        filter_note="Mevcut uygulama toplamsal puanlama kullanır: açık sözcüksel işaretler, embedding benzerliği ve öğrenilmiş geri besleme birlikte katkı verir.",
        feedback_bullets=[
            "Her kullanıcı için görünen son etiket",
            "Pozitif, negatif ve çelişkili kümeler",
            "Güncellenen sorgu ve çapa ifadeleri",
            "Güncellenen kavram ve kaynak öncelikleri",
        ],
        feedback_note="Mevcut durum: çalışan bir toplu geri besleme döngüsü vardır. Sonraki çalıştırmaları değiştirir, ancak henüz eğitilmiş bir sınıflandırıcı değildir.",
        current_now=[
            "Europe PMC, OpenAlex, Semantic Scholar ve DergiPark üzerinde çok kaynaklı arama",
            "PDF indirmeden önce iki aşamalı filtreleme",
            "Embedding destekli meta veri puanlama",
            "Geri besleme tabanlı istatistiksel terim güncellemesi",
            "Operasyonel döngüyü kapatan makale stoğu yenileme süreci",
        ],
        next_phase=[
            "Belge segmentasyonu",
            "LLM destekli çıkarım",
            "Yeterli etiket hacminden sonra sınıflandırıcı eğitimi",
        ],
        closing_statement="Arasınav için temel iddia şudur: erişim, anotasyon ve geri besleme artık tek bir sistem olarak birlikte çalışmaktadır.",
        notes_title="OpenNutri Yapay Zeka / Algoritma Mini Sunum Notları",
        speaker_notes=[
            "Slayt 1: Kendi rolünü erişim, sıralama, iki dilli arama ve geri besleme kullanımı olarak çerçevele. Giriş 30 saniyeyi geçmesin.",
            "Slayt 2: Problemin yalnızca scraping olmadığını vurgula. Asıl mesele hangi makalelerin indirmeye ve uzman zamanına değer olduğunu seçmektir.",
            "Slayt 3: Aşamalı hattı yavaş anlat. Buradaki en güçlü ifade 'Arama -> Filtre -> Edinim' olmalı.",
            "Slayt 4: Geri besleme döngüsünün uygulanmış olduğunu ama çevrim içi değil toplu güncellendiğini açıkça söyle. Bu en önemli ayrımlardan biridir.",
            "Slayt 5: Sınırları net çiz. Şu an çalışanları söyle, sonra ikinci döneme ertelenenleri savunmacı görünmeden belirt.",
        ],
    ),
}


def rgb_to_str(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide, color: RGBColor = GREEN) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_footer(slide, lang: str, slide_number: int) -> None:
    left = add_textbox(
        slide,
        0.55,
        7.03,
        4.2,
        0.22,
        "OpenNutri | AI / Algorithm mini deck"
        if lang == "en"
        else "OpenNutri | Yapay zeka / algoritma mini sunumu",
        10,
        MUTED,
    )
    left.paragraphs[0].alignment = PP_ALIGN.LEFT
    right = add_textbox(
        slide,
        11.7,
        7.03,
        1.0,
        0.22,
        f"{lang.upper()}  {slide_number:02d}",
        10,
        MUTED,
        bold=True,
    )
    right.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_panel(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor = PANEL) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    font_name: str = "Aptos",
) -> object:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return frame


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: Iterable[str],
    *,
    font_size: int = 18,
    bullet_color: RGBColor = INK,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.level = 0
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = bullet_color
        p.bullet = True
        p.space_after = Pt(7)


def add_chip(slide, left: float, top: float, width: float, text: str, fill_color: RGBColor, text_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, *, color: RGBColor = GREEN) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def build_title_slide(prs: Presentation, copy: DeckCopy) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.05),
        Inches(0.38),
        Inches(4.9),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    add_textbox(slide, 1.25, 1.0, 5.7, 1.35, copy.deck_title, 28, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, 1.28, 1.98, 5.8, 0.72, copy.deck_subtitle, 17, MUTED)

    add_panel(slide, 1.2, 2.58, 5.3, 2.12)
    add_textbox(slide, 1.5, 2.84, 1.9, 0.35, "Current focus" if copy.lang == "en" else "Mevcut odak", 13, GREEN, bold=True)
    add_bullets(slide, 1.48, 3.16, 4.7, 1.45, copy.title_scope, font_size=17)

    add_panel(slide, 7.2, 1.05, 5.4, 5.4, fill_color=GREEN_LIGHT)
    add_textbox(
        slide,
        7.55,
        1.35,
        4.6,
        0.5,
        "System claim at midterm" if copy.lang == "en" else "Arasınav sistem iddiası",
        15,
        GREEN,
        bold=True,
    )
    quote = (
        "Retrieval, annotation, and feedback already work together as one operational loop."
        if copy.lang == "en"
        else "Erişim, anotasyon ve geri besleme artık tek bir operasyonel döngü olarak birlikte çalışmaktadır."
    )
    add_textbox(slide, 7.55, 1.95, 4.35, 1.4, quote, 24, INK, bold=True, font_name="Aptos Display")
    add_textbox(
        slide,
        7.55,
        3.95,
        4.35,
        1.5,
        "Search, filtering, PDF acquisition, annotation, and later feedback reuse are all part of the current implementation."
        if copy.lang == "en"
        else "Arama, filtreleme, PDF edinimi, anotasyon ve sonraki geri besleme kullanımı mevcut uygulamanın parçasıdır.",
        17,
        MUTED,
    )

    add_footer(slide, copy.lang, 1)


def build_problem_slide(prs: Presentation, copy: DeckCopy) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)

    add_textbox(slide, 0.75, 0.55, 6.5, 0.6, copy.slide_titles[1], 24, INK, bold=True, font_name="Aptos Display")
    add_textbox(
        slide,
        0.78,
        1.05,
        10.8,
        0.45,
        "The algorithmic layer exists to protect expert time and improve retrieval quality."
        if copy.lang == "en"
        else "Algoritmik katman, uzman zamanını korumak ve erişim kalitesini yükseltmek için vardır.",
        16,
        MUTED,
    )

    card_lefts = [0.78, 4.45, 8.12]
    for idx, (title, body) in enumerate(copy.problem_cards):
        add_panel(slide, card_lefts[idx], 1.75, 3.25, 4.55)
        add_chip(slide, card_lefts[idx] + 0.25, 2.0, 1.55, f"0{idx + 1}", AMBER_LIGHT, AMBER)
        add_textbox(slide, card_lefts[idx] + 0.25, 2.55, 2.65, 1.05, title, 18, INK, bold=True, font_name="Aptos Display")
        add_textbox(slide, card_lefts[idx] + 0.25, 3.7, 2.7, 2.0, body, 16, MUTED)

    add_panel(slide, 0.78, 6.45, 10.6, 0.5, fill_color=AMBER_LIGHT)
    add_textbox(
        slide,
        1.05,
        6.59,
        10.1,
        0.25,
        "Goal: spend acquisition and annotation effort only on papers with real food-composition potential."
        if copy.lang == "en"
        else "Amaç: edinim ve anotasyon emeğini yalnızca gerçek gıda bileşimi potansiyeli taşıyan makalelere ayırmak.",
        14,
        INK,
        bold=True,
    )
    add_footer(slide, copy.lang, 2)


def build_pipeline_slide(prs: Presentation, copy: DeckCopy) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)

    add_textbox(slide, 0.75, 0.55, 7.5, 0.6, copy.slide_titles[2], 24, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, 0.78, 1.02, 10.9, 0.35, copy.pipeline_caption, 16, MUTED)

    left = 0.82
    top = 2.0
    widths = [2.05, 1.7, 2.05, 1.7, 2.0]
    colors = [GREEN_LIGHT, PANEL, AMBER_LIGHT, PANEL, GREEN_LIGHT]
    x = left
    for idx, label in enumerate(copy.pipeline_labels):
        add_panel(slide, x, top, widths[idx], 1.18, fill_color=colors[idx])
        add_textbox(slide, x + 0.18, top + 0.26, widths[idx] - 0.36, 0.62, label, 15, INK, bold=True, font_name="Aptos Display")
        if idx < len(copy.pipeline_labels) - 1:
            add_connector(slide, x + widths[idx], top + 0.59, x + widths[idx] + 0.25, top + 0.59)
        x += widths[idx] + 0.25

    add_textbox(
        slide,
        0.78,
        3.62,
        2.5,
        0.3,
        "Signals inside metadata filter" if copy.lang == "en" else "Meta veri filtresindeki sinyaller",
        13,
        GREEN,
        bold=True,
    )
    chip_positions = [
        (0.8, 4.05, 1.7),
        (2.65, 4.05, 1.95),
        (4.75, 4.05, 1.45),
        (6.35, 4.05, 1.75),
        (8.25, 4.05, 2.0),
        (3.7, 4.58, 2.25),
    ]
    for (x_pos, y_pos, width), text in zip(chip_positions, copy.filter_chips):
        fill = GREEN_LIGHT if x_pos < 6.0 else AMBER_LIGHT
        text_color = GREEN if x_pos < 6.0 else AMBER
        add_chip(slide, x_pos, y_pos, width, text, fill, text_color)

    add_panel(slide, 0.8, 5.28, 11.1, 0.92)
    add_textbox(slide, 1.05, 5.51, 10.6, 0.42, copy.filter_note, 15, MUTED)
    add_footer(slide, copy.lang, 3)


def build_feedback_slide(prs: Presentation, copy: DeckCopy) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)

    add_textbox(slide, 0.75, 0.55, 6.8, 0.6, copy.slide_titles[3], 24, INK, bold=True, font_name="Aptos Display")
    add_textbox(
        slide,
        0.78,
        1.02,
        10.8,
        0.35,
        "This loop is implemented now; it updates later crawler runs rather than only storing labels."
        if copy.lang == "en"
        else "Bu döngü şu anda uygulanmıştır; yalnızca etiket toplamaz, sonraki crawler çalıştırmalarını da günceller.",
        16,
        MUTED,
    )

    loop_boxes = [
        (1.1, 2.3, 2.25, 0.95, "Crawler run" if copy.lang == "en" else "Crawler çalıştırması"),
        (3.95, 1.58, 2.7, 0.95, "Search outcomes" if copy.lang == "en" else "Arama çıktıları"),
        (7.25, 2.3, 2.55, 0.95, "Annotation decisions" if copy.lang == "en" else "Anotasyon kararları"),
        (4.1, 4.45, 2.7, 0.95, "Feedback update" if copy.lang == "en" else "Geri besleme güncellemesi"),
    ]
    for x, y, w, h, label in loop_boxes:
        add_panel(slide, x, y, w, h, fill_color=GREEN_LIGHT if y < 3.0 else AMBER_LIGHT)
        add_textbox(slide, x + 0.18, y + 0.25, w - 0.36, 0.45, label, 18, INK, bold=True, font_name="Aptos Display")

    add_connector(slide, 3.35, 2.77, 3.95, 2.05)
    add_connector(slide, 6.65, 2.05, 7.25, 2.77)
    add_connector(slide, 8.45, 3.25, 6.05, 4.45)
    add_connector(slide, 4.75, 4.45, 2.2, 3.25)

    add_panel(slide, 10.25, 1.7, 2.35, 4.55)
    add_textbox(
        slide,
        10.55,
        1.98,
        1.75,
        0.3,
        "Update targets" if copy.lang == "en" else "Güncellenenler",
        13,
        GREEN,
        bold=True,
    )
    add_bullets(slide, 10.5, 2.35, 1.82, 2.85, copy.feedback_bullets, font_size=13)

    add_panel(slide, 0.8, 6.3, 11.75, 0.6, fill_color=AMBER_LIGHT)
    add_textbox(slide, 1.02, 6.48, 11.25, 0.25, copy.feedback_note, 12, INK, bold=True)
    add_footer(slide, copy.lang, 4)


def build_scope_slide(prs: Presentation, copy: DeckCopy) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)

    add_textbox(slide, 0.75, 0.55, 7.0, 0.6, copy.slide_titles[4], 24, INK, bold=True, font_name="Aptos Display")
    add_textbox(
        slide,
        0.78,
        1.05,
        10.9,
        0.35,
        "This slide helps you defend the project without overclaiming."
        if copy.lang == "en"
        else "Bu slayt, projeyi aşırı iddiaya kaçmadan savunmanıza yardımcı olur.",
        16,
        MUTED,
    )

    add_panel(slide, 0.82, 1.75, 5.75, 4.9, fill_color=GREEN_LIGHT)
    add_textbox(
        slide,
        1.08,
        2.0,
        2.7,
        0.55,
        "Implemented now" if copy.lang == "en" else "Şu anda uygulananlar",
        18,
        GREEN,
        bold=True,
        font_name="Aptos Display",
    )
    add_bullets(slide, 1.05, 2.62, 5.0, 3.35, copy.current_now, font_size=15)

    add_panel(slide, 6.85, 1.75, 5.65, 4.9, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.12, 2.0, 3.15, 0.65, "Deferred to second semester" if copy.lang == "en" else "İkinci döneme ertelenenler", 18, AMBER, bold=True, font_name="Aptos Display")
    add_bullets(slide, 7.08, 2.68, 4.85, 1.75, copy.next_phase, font_size=15)
    add_textbox(slide, 7.12, 4.6, 4.75, 1.25, copy.closing_statement, 17, INK, bold=True, font_name="Aptos Display")
    add_footer(slide, copy.lang, 5)


def build_notes(copy: DeckCopy, notes_path: Path) -> None:
    lines = [f"# {copy.notes_title}", ""]
    for idx, note in enumerate(copy.speaker_notes, start=1):
        lines.append(f"## Slide {idx}" if copy.lang == "en" else f"## Slayt {idx}")
        lines.append(note)
        lines.append("")
    notes_path.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")


def export_pdf(pptx_path: Path, pdf_dir: Path) -> None:
    pdf_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(pptx_path),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def build_deck(copy: DeckCopy) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_title_slide(prs, copy)
    build_problem_slide(prs, copy)
    build_pipeline_slide(prs, copy)
    build_feedback_slide(prs, copy)
    build_scope_slide(prs, copy)

    pptx_path = PPTX_DIR / (
        "OpenNutri_AI_Algorithm_Mini_Deck_EN.pptx"
        if copy.lang == "en"
        else "OpenNutri_AI_Algorithm_Mini_Deck_TR.pptx"
    )
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)

    notes_path = NOTES_DIR / (
        "OpenNutri_AI_Algorithm_Mini_Deck_EN.md"
        if copy.lang == "en"
        else "OpenNutri_AI_Algorithm_Mini_Deck_TR.md"
    )
    notes_path.parent.mkdir(parents=True, exist_ok=True)
    build_notes(copy, notes_path)
    return pptx_path


def main() -> None:
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    PDF_DIR.mkdir(parents=True, exist_ok=True)
    NOTES_DIR.mkdir(parents=True, exist_ok=True)

    for lang in ("en", "tr"):
        pptx_path = build_deck(COPY_BY_LANG[lang])
        export_pdf(pptx_path, PDF_DIR)
        print(f"Created {pptx_path}")

    print(f"Created notes in {NOTES_DIR}")


if __name__ == "__main__":
    main()
