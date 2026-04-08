from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx.dml.color import RGBColor

from export_ai_algorithm_slides import (
    AMBER,
    AMBER_LIGHT,
    BASE_DIR,
    GREEN,
    GREEN_LIGHT,
    INK,
    LINE,
    MUTED,
    NOTES_DIR,
    PANEL,
    PDF_DIR,
    PPTX_DIR,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    add_bullets,
    add_connector,
    add_panel,
    add_textbox,
    add_top_bar,
    export_pdf,
    set_background,
)


ASSETS_DIR = BASE_DIR / "assets"
BLUE = RGBColor(71, 115, 180)
BLUE_LIGHT = RGBColor(223, 233, 247)
RED = RGBColor(165, 75, 75)
RED_LIGHT = RGBColor(245, 225, 225)
PANEL_ALT = RGBColor(241, 236, 227)


def tr(lang: str, en: str, tur: str) -> str:
    return en if lang == "en" else tur


def asset(lang: str, stem: str) -> Path:
    suffix = "_en.png" if lang == "en" else ".png"
    return ASSETS_DIR / f"{stem}{suffix}"


def add_picture_panel(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    add_panel(slide, left, top, width, height)
    slide.shapes.add_picture(str(path), Inches(left + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(height - 0.16))


def add_small_card(slide, left: float, top: float, width: float, height: float, title: str, body: str, *, fill_color=PANEL, tag: str | None = None) -> None:
    add_panel(slide, left, top, width, height, fill_color=fill_color)
    if tag:
        tag_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left + 0.18),
            Inches(top + 0.16),
            Inches(0.9),
            Inches(0.34),
        )
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = AMBER_LIGHT
        tag_shape.line.fill.background()
        tag_frame = tag_shape.text_frame
        tag_frame.clear()
        p = tag_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag
        run.font.name = "Aptos"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = AMBER
        title_top = top + 0.62
    else:
        title_top = top + 0.22
    add_textbox(slide, left + 0.2, title_top, width - 0.4, 0.34, title, 15, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, left + 0.2, title_top + 0.44, width - 0.4, height - 0.75, body, 12, MUTED)


def add_hbar(slide, label: str, value: int, max_value: int, left: float, top: float, width: float, color: RGBColor) -> None:
    add_textbox(slide, left, top - 0.03, 1.9, 0.22, label, 11, INK)
    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left + 1.95),
        Inches(top),
        Inches(width),
        Inches(0.26),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = PANEL_ALT
    rail.line.fill.background()
    bar_width = width * (value / max_value) if max_value > 0 else 0.0
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left + 1.95),
        Inches(top),
        Inches(max(bar_width, 0.12)),
        Inches(0.26),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, left + 1.95 + width + 0.08, top - 0.04, 0.8, 0.22, str(value), 11, INK, bold=True)


def add_master_footer(slide, lang: str, slide_number: int) -> None:
    left = add_textbox(
        slide,
        0.55,
        7.02,
        5.2,
        0.22,
        tr(lang, "OpenNutri | AI / algorithm master deck", "OpenNutri | Yapay zeka / algoritma ana sunumu"),
        9,
        MUTED,
    )
    left.paragraphs[0].alignment = PP_ALIGN.LEFT
    right = add_textbox(slide, 11.65, 7.02, 1.0, 0.22, f"{lang.upper()}  {slide_number:02d}", 10, MUTED, bold=True)
    right.paragraphs[0].alignment = PP_ALIGN.RIGHT


def new_slide(prs: Presentation, lang: str, slide_number: int, title: str, subtitle: str | None = None, *, appendix: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, color=AMBER if appendix else GREEN)
    add_master_footer(slide, lang, slide_number)
    add_textbox(slide, 0.72, 0.45, 10.7, 0.5, title, 23, INK, bold=True, font_name="Aptos Display")
    if subtitle:
        add_textbox(slide, 0.75, 0.94, 11.6, 0.34, subtitle, 13, MUTED)
    return slide


def write_notes(path: Path, deck_title: str, slide_notes: list[tuple[int, str, str]], *, lang: str) -> None:
    lines = [f"# {deck_title}", ""]
    for number, title, body in slide_notes:
        lines.append(f"## Slide {number:02d} — {title}" if lang == "en" else f"## Slayt {number:02d} — {title}")
        lines.append(body.strip())
        lines.append("")
    path.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")


def load_run_summary() -> dict:
    manifest = (
        BASE_DIR.parent.parent
        / "services"
        / "data-pipeline"
        / "data"
        / "crawl_tr_live_2026-03-30b"
        / "raw_pdfs"
        / "_harvest_metadata.json"
    )
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    tr_summary = payload["summary"]["languages"]["tr"]
    terminal_recorded = (
        tr_summary["pdf_fetch_fail"]
        + tr_summary["pdf_validation_fail"]
        + tr_summary["accepted"]
    )
    return {
        **tr_summary,
        "terminal_recorded": terminal_recorded,
        "unaccounted": tr_summary["metadata_pass"] - terminal_recorded,
    }


def build_master_deck(lang: str) -> tuple[Path, Path]:
    run_summary = load_run_summary()
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    notes: list[tuple[int, str, str]] = []
    slide_no = 0

    def note(title: str, body_en: str, body_tr: str) -> None:
        notes.append((slide_no, title, tr(lang, body_en, body_tr)))

    deck_title = tr(lang, "OpenNutri AI / Algorithm Master Defense Notes", "OpenNutri Yapay Zeka / Algoritma Ana Savunma Notları")

    # 01. Title
    slide_no += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide)
    add_master_footer(slide, lang, slide_no)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(0.95),
        Inches(0.36),
        Inches(5.3),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    add_textbox(slide, 1.23, 0.92, 6.2, 1.1, tr(lang, "OpenNutri AI / Algorithm Master Defense", "OpenNutri Yapay Zeka / Algoritma Ana Savunması"), 28, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, 1.25, 1.92, 6.2, 0.52, tr(lang, "Midterm exposition deck | Retrieval, ranking, bilingual search, and feedback reuse", "Arasınav sunumu | Erişim, sıralama, iki dilli arama ve geri besleme kullanımı"), 16, MUTED)
    add_panel(slide, 1.18, 2.58, 5.35, 2.4)
    add_textbox(slide, 1.48, 2.86, 2.0, 0.3, tr(lang, "Role framing", "Rol çerçevesi"), 13, GREEN, bold=True)
    add_bullets(slide, 1.46, 3.16, 4.65, 1.65, [
        tr(lang, "Retrieval design across multiple scientific sources", "Birden fazla bilimsel kaynak üzerinde erişim tasarımı"),
        tr(lang, "Ranking logic before PDF acquisition", "PDF ediniminden önce sıralama mantığı"),
        tr(lang, "EN/TR language split and Turkish-source support", "EN/TR dil ayrımı ve Türkçe kaynak desteği"),
        tr(lang, "Feedback reuse from annotation outcomes", "Anotasyon çıktılarından geri besleme kullanımı"),
    ], font_size=15)
    add_panel(slide, 7.05, 1.05, 5.55, 5.45, fill_color=GREEN_LIGHT)
    add_textbox(slide, 7.42, 1.32, 3.0, 0.3, tr(lang, "Main midterm claim", "Arasınav ana iddiası"), 14, GREEN, bold=True)
    add_textbox(slide, 7.42, 1.88, 4.65, 1.9, tr(lang, "OpenNutri already operates as a closed loop linking search, filtering, PDF acquisition, annotation, and feedback reuse.", "OpenNutri şimdiden arama, filtreleme, PDF edinimi, anotasyon ve geri besleme kullanımını bağlayan kapalı bir döngü olarak çalışmaktadır."), 20, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, 7.42, 4.18, 4.75, 1.45, tr(lang, "This deck explains that loop from the AI / algorithm side, then goes deeper into the implementation details likely to matter in the defense.", "Bu sunum, bu döngüyü yapay zeka / algoritma tarafından açıklar ve ardından savunmada önemli olabilecek uygulama ayrıntılarına iner."), 15, MUTED)
    note(
        tr(lang, "Title and role framing", "Başlık ve rol çerçevesi"),
        "Start with the system claim, not with models. Say that the project already runs as one loop: search, filtering, acquisition, annotation, and feedback reuse.",
        "Modellerle değil, sistem iddiasıyla başla. Projenin şimdiden tek döngü olarak çalıştığını söyle: arama, filtreleme, edinim, anotasyon ve geri besleme kullanımı.",
    )

    # 02. Problem
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Why this layer exists", "Bu katman neden var"), tr(lang, "The algorithmic problem is not merely collecting papers; it is deciding which papers deserve download and expert time.", "Algoritmik problem yalnızca makale toplamak değildir; hangi makalelerin indirilmeye ve uzman zamanına değdiğini seçmektir."))
    add_small_card(slide, 0.78, 1.6, 3.27, 4.78, tr(lang, "Distributed literature", "Dağınık literatür"), tr(lang, "Relevant food-composition papers are spread across Europe PMC, OpenAlex, Semantic Scholar, DergiPark, and local open-access paths. A single-source strategy misses too much.", "Uygun gıda bileşimi makaleleri Europe PMC, OpenAlex, Semantic Scholar, DergiPark ve yerel açık erişim yollarına dağılmıştır. Tek kaynaklı strateji çok şey kaçırır."), tag="01")
    add_small_card(slide, 4.38, 1.6, 3.27, 4.78, tr(lang, "High acquisition cost", "Yüksek edinim maliyeti"), tr(lang, "Blindly downloading every candidate PDF is expensive, noisy, and wasteful. The system needs ranking before acquisition, not only after download.", "Her aday PDF'yi körlemesine indirmek pahalı, gürültülü ve verimsizdir. Sistem yalnızca indirmeden sonra değil, edinimden önce de sıralama gerektirir."), tag="02")
    add_small_card(slide, 7.98, 1.6, 3.27, 4.78, tr(lang, "Turkish visibility gap", "Türkçe görünürlük açığı"), tr(lang, "Turkish studies are underrepresented in standard international food-data systems. Retrieval therefore needs explicit EN/TR separation and Turkish-source support.", "Türkçe çalışmalar standart uluslararası gıda-veri sistemlerinde eksik temsil edilir. Bu nedenle erişim, açık EN/TR ayrımı ve Türkçe kaynak desteği gerektirir."), tag="03")
    add_panel(slide, 0.8, 6.5, 11.52, 0.42, fill_color=AMBER_LIGHT)
    add_textbox(slide, 1.02, 6.6, 11.0, 0.22, tr(lang, "Goal of the AI layer: maximize useful paper flow to the annotator while preserving recall and avoiding unnecessary PDF work.", "Yapay zeka katmanının amacı: geri çağırmayı korurken ve gereksiz PDF işini önlerken annotator'a giden yararlı makale akışını en üst düzeye çıkarmaktır."), 12, INK, bold=True)
    note(
        tr(lang, "Problem definition", "Problemin tanımı"),
        "Three pressures define the problem: the literature is distributed, acquisition is expensive, and Turkish coverage is weaker. This is why retrieval quality matters before annotation.",
        "Problemi üç baskı tanımlar: literatür dağınıktır, edinim pahalıdır ve Türkçe kapsam daha zayıftır. Bu yüzden erişim kalitesi anotasyondan önce önemlidir.",
    )

    # 03. System identity
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "What OpenNutri is at midterm", "OpenNutri arasınavda nedir"), tr(lang, "For the AI role, the important point is that the project is one connected loop rather than a separate scraper, backend, and UI.", "Yapay zeka rolü açısından önemli nokta, projenin ayrı bir scraper, backend ve arayüz değil, birbirine bağlı tek bir döngü olmasıdır."))
    labels = [tr(lang, "Search sources", "Arama kaynakları"), tr(lang, "Search gate", "Arama kapısı"), tr(lang, "Metadata filter", "Meta veri filtresi"), tr(lang, "PDF intake", "PDF edinimi"), tr(lang, "Expert annotation", "Uzman anotasyonu"), tr(lang, "Feedback update", "Geri besleme güncellemesi")]
    xs = [0.85, 2.55, 4.35, 6.25, 8.15, 10.05]
    widths = [1.45, 1.45, 1.55, 1.45, 1.65, 1.75]
    fills = [GREEN_LIGHT, PANEL, PANEL, PANEL, GREEN_LIGHT, AMBER_LIGHT]
    for idx, (x, w, label, fill) in enumerate(zip(xs, widths, labels, fills)):
        add_panel(slide, x, 2.0, w, 1.08, fill_color=fill)
        box = add_textbox(slide, x + 0.1, 2.28, w - 0.2, 0.38, label, 12, INK, bold=True, font_name="Aptos Display")
        box.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx < len(xs) - 1:
            add_connector(slide, x + w, 2.54, xs[idx + 1], 2.54)
    add_panel(slide, 0.9, 3.75, 5.7, 2.55, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.2, 4.02, 3.2, 0.3, tr(lang, "What it is", "Ne olduğudur"), 13, GREEN, bold=True)
    add_bullets(slide, 1.18, 4.34, 4.95, 1.72, [
        tr(lang, "A human-in-the-loop scientific paper discovery system", "İnsan döngüsünde çalışan bilimsel makale keşif sistemi"),
        tr(lang, "A shared data contract between crawler, storage, and annotation", "Crawler, depolama ve anotasyon arasında ortak veri sözleşmesi"),
        tr(lang, "A feedback-driven retrieval system that changes later runs", "Sonraki çalıştırmaları değiştiren geri besleme güdümlü erişim sistemi"),
    ], font_size=14)
    add_panel(slide, 6.82, 3.75, 5.6, 2.55, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.1, 4.02, 3.2, 0.3, tr(lang, "What it is not", "Ne olmadığıdır"), 13, AMBER, bold=True)
    add_bullets(slide, 7.08, 4.34, 4.88, 1.72, [
        tr(lang, "Not a one-shot scraper that downloads everything", "Her şeyi indiren tek seferlik bir scraper değildir"),
        tr(lang, "Not only a database schema or only a labeling UI", "Yalnızca bir veritabanı şeması veya yalnızca bir etiketleme arayüzü değildir"),
        tr(lang, "Not yet a fully automatic extraction engine", "Henüz tamamen otomatik bir çıkarım motoru değildir"),
    ], font_size=14)
    note(
        tr(lang, "System identity", "Sistemin kimliği"),
        "Use this slide to say what you are presenting and what you are not. The best line is: we are presenting a retrieval system that feeds and learns from expert annotation.",
        "Bu slaytı ne sunduğunuzu ve ne sunmadığınızı söylemek için kullanın. En iyi cümle: uzman anotasyonundan beslenen ve ondan öğrenen bir erişim sistemi sunuyoruz.",
    )

    # 04. Architecture
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "End-to-end architecture", "Uçtan uca mimari"), tr(lang, "The architecture already links discovery, storage, annotation, and later feedback reuse. This is the system map for the rest of the presentation.", "Mimari, keşif, depolama, anotasyon ve sonraki geri besleme kullanımını şimdiden bağlıyor. Bu, sunumun geri kalanı için sistem haritasıdır."))
    add_picture_panel(slide, asset(lang, "figure_1_system_architecture"), 6.0, 1.55, 6.35, 5.32)
    add_panel(slide, 0.82, 1.55, 4.82, 5.32)
    add_textbox(slide, 1.1, 1.84, 2.8, 0.3, tr(lang, "How to read this figure", "Bu şekil nasıl okunur"), 14, GREEN, bold=True)
    add_bullets(slide, 1.08, 2.24, 4.1, 3.42, [
        tr(lang, "Search sources produce metadata candidates.", "Arama kaynakları meta veri adayları üretir."),
        tr(lang, "The crawler decides which candidates deserve acquisition.", "Crawler hangi adayların edinime değdiğine karar verir."),
        tr(lang, "Accepted papers move into Supabase storage and paper records.", "Kabul edilen makaleler Supabase depolamasına ve paper kayıtlarına taşınır."),
        tr(lang, "The annotator creates structured food-item and nutrient evidence.", "Annotator yapılandırılmış food-item ve nutrient kanıtı oluşturur."),
        tr(lang, "Label events and global labels later feed the next retrieval cycle.", "Label event'ler ve global label'lar sonraki erişim döngüsünü besler."),
    ], font_size=13)
    add_panel(slide, 0.95, 6.02, 4.45, 0.5, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.12, 6.17, 4.0, 0.18, tr(lang, "Main defense point: the AI layer is connected to real user decisions, not isolated from the product.", "Ana savunma noktası: yapay zeka katmanı üründen kopuk değil, gerçek kullanıcı kararlarına bağlıdır."), 11, INK, bold=True)
    note(
        tr(lang, "Architecture map", "Mimari haritası"),
        "Walk this figure slowly from left to right and then close the loop. Refer back to it when later slides get detailed.",
        "Bu şekli soldan sağa yavaş anlatın ve sonra döngüyü kapatın. Sonraki slaytlar ayrıntılandığında buna geri dönün.",
    )

    # 05. Data model
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Shared data model relevant to AI", "Yapay zeka açısından ilgili ortak veri modeli"), tr(lang, "The retrieval layer depends on papers, search evidence, user events, and a shared vocabulary.", "Erişim katmanı paper kayıtlarına, arama kanıtına, kullanıcı event'lerine ve ortak bir sözlüğe dayanır."))
    add_picture_panel(slide, asset(lang, "figure_3_database_schema"), 6.02, 1.55, 6.28, 5.28)
    add_panel(slide, 0.82, 1.55, 4.82, 5.28)
    add_bullets(slide, 1.08, 1.86, 4.1, 4.18, [
        tr(lang, "Reference vocabulary: entities, aliases, and master nutrients are reused by both retrieval and annotation.", "Referans sözlük: entities, alias'lar ve master nutrient'ler hem erişim hem anotasyon tarafından yeniden kullanılır."),
        tr(lang, "Paper intake: papers store accepted PDF-level records available to the UI.", "Paper alımı: papers tablosu, arayüz için mevcut kabul edilmiş PDF düzeyi kayıtları tutar."),
        tr(lang, "Search evidence: paper_search_hits, paper_search_batches, and paper_search_batch_hits preserve why a paper entered the system.", "Arama kanıtı: paper_search_hits, paper_search_batches ve paper_search_batch_hits, bir paper'ın sisteme neden girdiğini korur."),
        tr(lang, "Feedback evidence: paper_label_events, paper_global_labels, and search_sessions preserve user behavior for later ranking updates.", "Geri besleme kanıtı: paper_label_events, paper_global_labels ve search_sessions, sonraki sıralama güncellemeleri için kullanıcı davranışını korur."),
    ], font_size=13)
    add_panel(slide, 0.95, 6.05, 4.45, 0.52, fill_color=AMBER_LIGHT)
    add_textbox(slide, 1.12, 6.2, 4.0, 0.2, tr(lang, "Important distinction: raw search evidence is stored separately from accepted paper records.", "Önemli ayrım: ham arama kanıtı, kabul edilmiş paper kayıtlarından ayrı tutulur."), 11, INK, bold=True)
    note(
        tr(lang, "Database slice for AI", "Yapay zeka için veritabanı kesiti"),
        "Emphasize that the AI layer does not only consume papers. It also depends on search-evidence tables and label-evidence tables.",
        "Yapay zeka katmanının yalnızca paper kayıtlarını kullanmadığını vurgulayın. Arama kanıtı tablolarına ve etiket kanıtı tablolarına da dayanır.",
    )

    # 06. Search layer
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Search layer and language split", "Arama katmanı ve dil ayrımı"), tr(lang, "The crawler is multi-source by design, and it treats English and Turkish as different retrieval environments rather than one merged pool.", "Crawler tasarım gereği çok kaynaklıdır ve İngilizce ile Türkçeyi tek birleşik havuz yerine farklı erişim ortamları olarak ele alır."))
    card_data = [
        ("Europe PMC", tr(lang, "Main open-access biomedical surface; strongest default bias for English.", "Ana açık erişimli biyomedikal yüzey; İngilizce için en güçlü varsayılan öncelik.")),
        ("OpenAlex", tr(lang, "Broad metadata coverage used in both languages.", "Her iki dilde de kullanılan geniş meta veri kapsamı.")),
        ("Semantic Scholar", tr(lang, "Additional discovery path and duplicate-aware evidence.", "Ek keşif yolu ve duplicate farkındalıklı kanıt.")),
        ("DergiPark", tr(lang, "Essential Turkish source; highest static bias in Turkish.", "Temel Türkçe kaynak; Türkçede en yüksek statik öncelik.")),
    ]
    for x, (title, body) in zip([0.78, 3.95, 7.12, 10.29], card_data):
        add_small_card(slide, x, 1.68, 2.28, 2.68, title, body, fill_color=GREEN_LIGHT if title == "DergiPark" else PANEL)
    add_panel(slide, 0.8, 4.72, 5.85, 1.83, fill_color=BLUE_LIGHT)
    add_textbox(slide, 1.08, 4.98, 2.4, 0.3, "EN", 16, BLUE, bold=True, font_name="Aptos Display")
    add_bullets(slide, 1.05, 5.32, 4.95, 1.0, [
        tr(lang, "Default bias favors Europe PMC, then OpenAlex.", "Varsayılan öncelik Europe PMC'yi, sonra OpenAlex'i öne çıkarır."),
        tr(lang, "English composition frames and phrase pools stay language-scoped.", "İngilizce composition frame'ler ve phrase havuzları dil kapsamlı kalır."),
    ], font_size=12)
    add_panel(slide, 6.85, 4.72, 5.55, 1.83, fill_color=GREEN_LIGHT)
    add_textbox(slide, 7.15, 4.98, 2.4, 0.3, "TR", 16, GREEN, bold=True, font_name="Aptos Display")
    add_bullets(slide, 7.12, 5.32, 4.75, 1.0, [
        tr(lang, "Default bias favors DergiPark, then OpenAlex and Semantic Scholar.", "Varsayılan öncelik DergiPark'ı, sonra OpenAlex ve Semantic Scholar'ı öne çıkarır."),
        tr(lang, "Separate anchors, phrases, and feedback scores keep Turkish behavior language-scoped.", "Ayrı anchor, phrase ve feedback score yapıları Türkçe davranışı dil kapsamlı tutar."),
    ], font_size=12)
    note(
        tr(lang, "Source strategy", "Kaynak stratejisi"),
        "Say explicitly that each source fills a different gap. The language split is not cosmetic; it changes source priors, query phrases, and search order.",
        "Her kaynağın farklı bir açığı kapattığını açıkça söyleyin. Dil ayrımı kozmetik değildir; source prior'ları, query phrase'leri ve arama sırasını değiştirir.",
    )

    # 07. Query construction
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Query construction", "Sorgu oluşturma"), tr(lang, "The crawler does not rely on a fixed query list. It mixes base composition queries, learned query phrases, and a rotating concept pool.", "Crawler sabit bir sorgu listesine dayanmaz. Temel composition sorgularını, öğrenilmiş query phrase'leri ve dönen bir concept havuzunu birleştirir."))
    labels = [tr(lang, "Base queries", "Temel sorgular"), tr(lang, "Concept pool", "Concept havuzu"), tr(lang, "Phrase pool", "Phrase havuzu"), tr(lang, "Learned query", "Öğrenilmiş sorgu"), tr(lang, "Source task", "Kaynak görevi")]
    xs = [0.82, 2.78, 4.74, 6.82, 9.02]
    widths = [1.56, 1.56, 1.56, 1.78, 2.18]
    fills = [BLUE_LIGHT, GREEN_LIGHT, GREEN_LIGHT, AMBER_LIGHT, PANEL]
    for idx, (x, w, label, fill) in enumerate(zip(xs, widths, labels, fills)):
        add_panel(slide, x, 1.78, w, 1.08, fill_color=fill)
        frame = add_textbox(slide, x + 0.12, 2.1, w - 0.24, 0.34, label, 12, INK, bold=True, font_name="Aptos Display")
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx < len(xs) - 1:
            add_connector(slide, x + w, 2.52, xs[idx + 1], 2.52)
    add_panel(slide, 0.82, 3.35, 5.9, 2.8)
    add_bullets(slide, 1.06, 3.72, 5.02, 1.92, [
        tr(lang, "Start with language-specific base composition templates.", "Dil-özgü temel composition şablonlarıyla başlar."),
        tr(lang, "Interleave food terms and nutrient terms in one concept pool.", "Food term'leri ve nutrient term'leri tek concept havuzunda iç içe geçirir."),
        tr(lang, "Reorder concepts using concept_scores learned from previous outcomes.", "Concept'leri önceki sonuçlardan öğrenilen concept_scores ile yeniden sıralar."),
        tr(lang, "Mix core phrases with a small exploration slice from the longer phrase pool.", "Çekirdek phrase'leri, daha uzun havuzdan küçük bir exploration bölümüyle karıştırır."),
    ], font_size=12)
    add_panel(slide, 7.02, 3.35, 5.38, 2.8, fill_color=GREEN_LIGHT)
    add_textbox(slide, 7.3, 3.62, 3.0, 0.3, tr(lang, "Exact implementation choices", "Tam uygulama tercihleri"), 14, GREEN, bold=True)
    add_bullets(slide, 7.28, 4.0, 4.65, 1.82, [
        tr(lang, "Exploration rate for extra phrases: 2 percent of the phrase pool.", "Ek phrase'ler için exploration oranı: phrase havuzunun yüzde 2'si."),
        tr(lang, "Separate cursors are preserved per language.", "Her dil için ayrı cursor korunur."),
        tr(lang, "Nutrient learned queries add unit framing such as mg/100g or g/100g.", "Nutrient odaklı öğrenilmiş sorgular mg/100g veya g/100g gibi birim çerçevesi ekler."),
        tr(lang, "Queries are normalized and deduplicated before search-task expansion.", "Sorgular arama görevi genişletmesinden önce normalize edilip eşsizleştirilir."),
    ], font_size=12)
    note(
        tr(lang, "Query generation logic", "Sorgu üretim mantığı"),
        "The committee should hear that queries are partly seeded and partly learned. The phrase pool keeps a tiny exploration slice so the system does not get trapped in one phrase set.",
        "Kurul, sorguların kısmen seed kısmen öğrenilmiş olduğunu duymalı. Phrase havuzu, sistemin tek bir phrase setine sıkışmaması için küçük bir exploration bölümü tutar.",
    )

    # 08. Prioritization
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Search-task prioritization", "Arama görevi önceliklendirmesi"), tr(lang, "After query construction, each query is expanded into source-specific search tasks and ranked before execution.", "Sorgu oluşturulduktan sonra her sorgu kaynağa özgü arama görevlerine genişletilir ve çalıştırmadan önce sıralanır."))
    add_panel(slide, 0.82, 1.62, 5.9, 5.0, fill_color=BLUE_LIGHT)
    add_textbox(slide, 1.1, 1.9, 3.0, 0.34, tr(lang, "Priority score", "Öncelik puanı"), 18, BLUE, bold=True, font_name="Aptos Display")
    add_textbox(slide, 1.12, 2.42, 5.0, 0.72, "priority_score = pair_score + batch_score", 24, INK, bold=True, font_name="Consolas")
    add_bullets(slide, 1.1, 3.3, 4.95, 2.1, [
        tr(lang, "pair_score = learned pair feedback + 0.15 x source_prior + static source bias", "pair_score = öğrenilmiş pair feedback + 0.15 x source_prior + statik kaynak önceliği"),
        tr(lang, "batch_score comes from historical performance of the exact batch key", "batch_score, tam batch key'in tarihsel performansından gelir"),
        tr(lang, "Tasks are sorted descending and truncated to the run budget", "Görevler azalan sırada sıralanır ve çalıştırma bütçesine göre kesilir"),
    ], font_size=13)
    add_panel(slide, 7.0, 1.62, 5.35, 5.0)
    add_textbox(slide, 7.3, 1.9, 3.4, 0.34, tr(lang, "Static source bias by language", "Dile göre statik kaynak önceliği"), 16, GREEN, bold=True, font_name="Aptos Display")
    add_textbox(slide, 7.35, 2.4, 4.6, 2.65, tr(lang, "EN:\nEurope PMC +0.35\nOpenAlex +0.15\nSemantic Scholar +0.05\nDergiPark -0.15\n\nTR:\nDergiPark +0.35\nOpenAlex +0.20\nSemantic Scholar +0.10\nEurope PMC -0.30", "EN:\nEurope PMC +0.35\nOpenAlex +0.15\nSemantic Scholar +0.05\nDergiPark -0.15\n\nTR:\nDergiPark +0.35\nOpenAlex +0.20\nSemantic Scholar +0.10\nEurope PMC -0.30"), 15, INK, bold=True, font_name="Consolas")
    add_panel(slide, 7.28, 5.4, 4.6, 0.72, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.5, 5.63, 4.15, 0.2, tr(lang, "This is how the system turns language strategy into execution order.", "Sistem dil stratejisini bu şekilde çalıştırma sırasına dönüştürür."), 12, INK, bold=True)
    note(
        tr(lang, "Task ranking", "Görev sıralaması"),
        "Pair scores capture source-template-term performance. Batch scores capture the exact query batch. Static bias is the explicit language-policy layer on top of learned feedback.",
        "Pair score'lar source-template-term performansını yakalar. Batch score'lar tam query batch'i yakalar. Statik öncelik, öğrenilmiş geri beslemenin üstündeki açık dil politikası katmanıdır.",
    )

    # 09. Search gate
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Search gate", "Arama kapısı"), tr(lang, "The search gate is a cheap first-pass metadata screen. It rejects weak candidates before the heavier metadata filter and acquisition stages.", "Arama kapısı ucuz bir ilk meta veri taramasıdır. Daha ağır meta veri filtresi ve edinim aşamalarından önce zayıf adayları reddeder."))
    add_panel(slide, 0.82, 1.62, 5.6, 4.95, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.1, 1.92, 2.8, 0.32, tr(lang, "Positive signals", "Pozitif sinyaller"), 16, GREEN, bold=True, font_name="Aptos Display")
    add_bullets(slide, 1.08, 2.35, 4.8, 3.65, [
        tr(lang, "First anchor/composition phrase hit (+0.90)", "İlk anchor/composition phrase eşleşmesi (+0.90)"),
        tr(lang, "First food term hit (+0.35)", "İlk food term eşleşmesi (+0.35)"),
        tr(lang, "First nutrient term hit (+0.35)", "İlk nutrient term eşleşmesi (+0.35)"),
        tr(lang, "Unit pattern such as mg/100g or g/100g (+0.70)", "mg/100g veya g/100g gibi birim örüntüsü (+0.70)"),
        tr(lang, "Food + nutrient combo bonus (+0.45)", "Food + nutrient kombinasyon bonusu (+0.45)"),
    ], font_size=14)
    add_panel(slide, 6.7, 1.62, 5.6, 4.95, fill_color=RED_LIGHT)
    add_textbox(slide, 6.98, 1.92, 3.1, 0.32, tr(lang, "Penalties and pass rule", "Cezalar ve geçiş kuralı"), 16, RED, bold=True, font_name="Aptos Display")
    add_bullets(slide, 6.98, 2.35, 4.82, 3.65, [
        tr(lang, "Missing metadata -> immediate reject with score -2.0", "Eksik meta veri -> -2.0 puanla doğrudan ret"),
        tr(lang, "Missing abstract penalty: -1.30 or -0.45", "Eksik özet cezası: -1.30 veya -0.45"),
        tr(lang, "Strong negative signals penalized up to -2.5", "Güçlü negatif sinyaller en fazla -2.5 ceza alır"),
        tr(lang, "Soft negatives penalized up to -1.5", "Yumuşak negatifler en fazla -1.5 ceza alır"),
        tr(lang, "Health-outcome terms penalized up to -1.6", "Sağlık sonucu terimleri en fazla -1.6 ceza alır"),
        tr(lang, "Pass threshold: score >= -0.35", "Geçiş eşiği: score >= -0.35"),
    ], font_size=14)
    note(
        tr(lang, "First-pass filter", "İlk geçiş filtresi"),
        "Stress that the search gate is intentionally cheap and tolerant. The threshold is negative because this gate is not the final decision layer.",
        "Arama kapısının bilinçli olarak ucuz ve toleranslı olduğunu vurgulayın. Eşik negatiftir çünkü bu kapı son karar katmanı değildir.",
    )

    # 10. Metadata filter
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Metadata filter", "Meta veri filtresi"), tr(lang, "The metadata filter is the stronger additive decision layer. It combines lexical cues, source priors, embeddings, and feedback terms before PDF acquisition.", "Meta veri filtresi daha güçlü toplamsal karar katmanıdır. PDF ediniminden önce sözcüksel işaretleri, kaynak önceliklerini, embedding'leri ve feedback term'lerini birleştirir."))
    add_panel(slide, 0.82, 1.62, 5.92, 4.95, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.1, 1.92, 3.2, 0.32, tr(lang, "Core evidence", "Temel kanıt"), 16, GREEN, bold=True, font_name="Aptos Display")
    add_bullets(slide, 1.08, 2.35, 5.1, 3.62, [
        tr(lang, "Composition phrase +1.35", "Composition phrase +1.35"),
        tr(lang, "Unit pattern +1.25", "Birim örüntüsü +1.25"),
        tr(lang, "Food hit +0.65 and nutrient hit +0.65", "Food hit +0.65 ve nutrient hit +0.65"),
        tr(lang, "Food + nutrient combo +0.75", "Food + nutrient kombinasyonu +0.75"),
        tr(lang, "Missing abstract penalty -1.10 or -0.35", "Eksik özet cezası -1.10 veya -0.35"),
        tr(lang, "Health-outcome penalty capped at -2.0", "Sağlık sonucu cezası en fazla -2.0"),
    ], font_size=14)
    add_panel(slide, 6.95, 1.62, 5.33, 4.95, fill_color=BLUE_LIGHT)
    add_textbox(slide, 7.22, 1.92, 3.2, 0.32, tr(lang, "Learned and semantic evidence", "Öğrenilmiş ve anlamsal kanıt"), 16, BLUE, bold=True, font_name="Aptos Display")
    add_bullets(slide, 7.2, 2.35, 4.58, 3.62, [
        tr(lang, "Source prior from historical positive/negative balance, clamped to +/-0.9", "Tarihsel pozitif/negatif dengesinden gelen kaynak önceliği, +/-0.9 ile sınırlandırılır"),
        tr(lang, "Embedding similarity adds +1.45 if similarity crosses the language threshold", "Embedding benzerliği, benzerlik dil eşiğini aşarsa +1.45 ekler"),
        tr(lang, "An extra +0.75 is added when embedding similarity is above threshold", "Embedding benzerliği eşik üzerindeyse ek +0.75 eklenir"),
        tr(lang, "Feedback n-gram score is summed and clamped to +/-6.0", "Feedback n-gram skoru toplanır ve +/-6.0 ile sınırlandırılır"),
        tr(lang, "Accept threshold: score >= 1.75", "Kabul eşiği: score >= 1.75"),
    ], font_size=13)
    note(
        tr(lang, "Main ranking layer", "Ana sıralama katmanı"),
        "This is the core ranking story: additive scoring from several evidence families. No single signal decides the paper alone.",
        "Sıralama hikâyesinin çekirdeği budur: birkaç kanıt ailesinden gelen toplamsal puanlama. Tek bir sinyal tek başına makaleye karar vermez.",
    )

    # 11. Embedding layer
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Embedding layer", "Embedding katmanı"), tr(lang, "Embeddings are a semantic support layer inside metadata scoring, not a standalone classifier.", "Embedding'ler bağımsız bir sınıflandırıcı değil, meta veri puanlamasının içindeki anlamsal destek katmanıdır."))
    add_small_card(slide, 0.82, 1.65, 3.45, 2.05, tr(lang, "English model", "İngilizce model"), "all-MiniLM-L6-v2\nthreshold = 0.45", fill_color=BLUE_LIGHT)
    add_small_card(slide, 4.55, 1.65, 3.55, 2.05, tr(lang, "Turkish / multilingual model", "Türkçe / çok dilli model"), "paraphrase-multilingual-MiniLM-L12-v2\nthreshold = 0.42", fill_color=GREEN_LIGHT)
    add_small_card(slide, 8.38, 1.65, 3.95, 2.05, tr(lang, "Shared rule", "Ortak kural"), tr(lang, "Trim metadata text to 1800 characters before encoding.", "Encoding'den önce meta veri metnini 1800 karaktere kırp."), fill_color=AMBER_LIGHT)
    add_panel(slide, 0.82, 4.02, 5.6, 2.55)
    add_textbox(slide, 1.1, 4.28, 3.0, 0.32, tr(lang, "How the score is produced", "Skor nasıl üretilir"), 15, GREEN, bold=True)
    add_bullets(slide, 1.08, 4.68, 4.85, 1.48, [
        tr(lang, "Load anchor phrases from feedback config; fall back to seeded anchors if needed.", "Anchor phrase'leri feedback config'den yükle; gerekirse seed anchor'lara dön."),
        tr(lang, "Encode the candidate title + abstract in the language-specific model.", "Aday başlık + özeti dil-özgü modelde encode et."),
        tr(lang, "Compute normalized similarity against all anchors and keep the maximum match.", "Tüm anchor'lara karşı normalize benzerliği hesapla ve en yüksek eşleşmeyi tut."),
    ], font_size=12)
    add_panel(slide, 6.7, 4.02, 5.6, 2.55, fill_color=AMBER_LIGHT)
    add_textbox(slide, 6.98, 4.28, 3.2, 0.32, tr(lang, "Why this design", "Bu tasarım neden seçildi"), 15, AMBER, bold=True)
    add_bullets(slide, 6.95, 4.68, 4.85, 1.48, [
        tr(lang, "Anchors keep the semantic target narrow: composition-oriented rather than general nutrition.", "Anchor'lar anlamsal hedefi dar tutar: genel beslenmeden çok composition odaklı."),
        tr(lang, "Two models keep EN and TR scoring language-aware.", "İki model EN ve TR puanlamasını dil duyarlı tutar."),
        tr(lang, "This layer recovers semantically relevant phrasing that exact keywords may miss.", "Bu katman, tam anahtar kelimelerin kaçırabileceği anlamsal olarak ilgili ifadeleri yakalar."),
    ], font_size=12)
    note(
        tr(lang, "Embedding design", "Embedding tasarımı"),
        "If asked which embeddings are used, answer exactly: English all-MiniLM-L6-v2, Turkish/multilingual paraphrase-multilingual-MiniLM-L12-v2. Then say embeddings support metadata scoring; they do not replace the rest of the ranking logic.",
        "Hangi embedding'lerin kullanıldığını sorarlarsa tam olarak şöyle cevap verin: İngilizce all-MiniLM-L6-v2, Türkçe/çok dilli paraphrase-multilingual-MiniLM-L12-v2. Sonra embedding'lerin meta veri puanlamasını desteklediğini, sıralamanın geri kalanını değiştirmediğini söyleyin.",
    )

    # 12. Feedback labels
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Feedback labels and training sets", "Geri besleme etiketleri ve eğitim kümeleri"), tr(lang, "The learning loop is based on the latest visible state of each annotator, not on raw event totals.", "Öğrenme döngüsü ham event toplamlarına değil, her annotator'ün görünen son durumuna dayanır."))
    add_picture_panel(slide, asset(lang, "figure_2_feedback_data_model"), 6.05, 1.55, 6.22, 5.28)
    add_panel(slide, 0.82, 1.55, 4.92, 5.28)
    add_bullets(slide, 1.08, 1.86, 4.15, 4.25, [
        tr(lang, "Positive paper: latest draft or done state with has_data = true, food-item count > 0, and nutrient-value count > 0.", "Pozitif paper: has_data = true, food-item sayısı > 0 ve nutrient-value sayısı > 0 olan son draft veya done durumu."),
        tr(lang, "Negative paper: a global definitely-no-data label or at least two unique skip users.", "Negatif paper: global definitely-no-data etiketi veya en az iki farklı skip kullanıcısı."),
        tr(lang, "Conflict paper: both positive and negative evidence; excluded from learning.", "Çelişkili paper: hem pozitif hem negatif kanıt içerir; öğrenmeden çıkarılır."),
        tr(lang, "Processing is language-scoped after label construction.", "İşleme, etiket kurulumu sonrası dil-özgü yürütülür."),
    ], font_size=13)
    add_panel(slide, 0.95, 6.02, 4.45, 0.5, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.12, 6.17, 4.0, 0.18, tr(lang, "Important: the feedback loop is implemented, but it is batch-updated rather than online.", "Önemli nokta: geri besleme döngüsü uygulanmıştır, ancak çevrim içi değil toplu güncellenir."), 11, INK, bold=True)
    note(
        tr(lang, "Label semantics", "Etiket semantiği"),
        "This slide answers what exactly becomes training data. Emphasize latest visible state, not raw counts, and say conflicts are excluded on purpose.",
        "Bu slayt tam olarak neyin eğitim verisine dönüştüğünü cevaplar. Ham sayılar yerine görünen son durumu vurgulayın ve çatışmaların bilinçli olarak dışlandığını söyleyin.",
    )

    # 13. Feedback updates
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "What feedback changes in later runs", "Geri besleme sonraki çalıştırmalarda neyi değiştirir"), tr(lang, "The loop is already active. Feedback is not only stored; it is exported into a new config and consumed by the next crawler run.", "Döngü hâlihazırda aktiftir. Geri besleme yalnızca saklanmaz; yeni bir config'e aktarılır ve bir sonraki crawler çalıştırması tarafından kullanılır."))
    boxes = [tr(lang, "Annotation outcomes", "Anotasyon çıktıları"), tr(lang, "Feedback export", "Feedback export"), "latest.json", tr(lang, "Crawler startup", "Crawler başlangıcı"), tr(lang, "Later search behavior", "Sonraki arama davranışı")]
    xs = [0.85, 2.9, 5.15, 7.35, 9.55]
    widths = [1.55, 1.75, 1.55, 1.65, 2.05]
    fills = [GREEN_LIGHT, AMBER_LIGHT, BLUE_LIGHT, PANEL, GREEN_LIGHT]
    for idx, (x, w, label, fill) in enumerate(zip(xs, widths, boxes, fills)):
        add_panel(slide, x, 2.08, w, 1.05, fill_color=fill)
        frame = add_textbox(slide, x + 0.12, 2.38, w - 0.24, 0.34, label, 12, INK, bold=True, font_name="Aptos Display")
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx < len(xs) - 1:
            add_connector(slide, x + w, 2.57, xs[idx + 1], 2.57)
    add_panel(slide, 0.85, 3.72, 5.8, 2.72)
    add_bullets(slide, 1.1, 4.08, 4.95, 1.8, [
        "query_phrases",
        "anchor_phrases",
        "weighted_terms",
        "concept_scores",
        "pair_scores, batch_scores, source_priors",
    ], font_size=14)
    add_panel(slide, 6.95, 3.72, 5.32, 2.72, fill_color=GREEN_LIGHT)
    add_bullets(slide, 7.2, 4.08, 4.6, 1.8, [
        tr(lang, "Later runs search different phrases in different orders.", "Sonraki çalıştırmalar farklı phrase'leri farklı sıralarda arar."),
        tr(lang, "Concept ordering changes which foods or nutrients are combined first.", "Concept sıralaması hangi food veya nutrient'lerin önce birleşeceğini değiştirir."),
        tr(lang, "Source-query combinations with better yield move up the queue.", "Daha iyi verim veren source-query kombinasyonları kuyruğun üstüne çıkar."),
    ], font_size=13)
    note(
        tr(lang, "Active feedback loop", "Aktif geri besleme döngüsü"),
        "If asked whether the loop is implemented or only collecting signals, answer: implemented, batch-updated, and already changing later runs.",
        "Döngünün uygulanmış mı yoksa sadece sinyal topluyor mu olduğunu sorarlarsa şöyle cevap verin: uygulanmış, toplu güncelleniyor ve sonraki çalıştırmaları şimdiden değiştiriyor.",
    )

    # 14. Operational closure
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Operational closure", "Operasyonel kapanış"), tr(lang, "The refill script is what turns the algorithmic pieces into a working operational loop when the annotator queue runs low.", "Refill betiği, annotator kuyruğu azaldığında algoritmik parçaları çalışan operasyonel bir döngüye dönüştüren bileşendir."))
    steps = [tr(lang, "Check EN/TR deficits", "EN/TR açıklarını kontrol et"), tr(lang, "Refresh feedback terms", "Feedback term'lerini yenile"), tr(lang, "Refresh DergiPark index", "DergiPark indeksini yenile"), tr(lang, "Run crawler v2", "crawler v2 çalıştır"), tr(lang, "Upload accepted PDFs", "Kabul edilen PDF'leri yükle")]
    x = 0.86
    for idx, step in enumerate(steps):
        add_panel(slide, x, 2.0, 2.08, 1.08, fill_color=GREEN_LIGHT if idx in {1, 3} else PANEL)
        frame = add_textbox(slide, x + 0.14, 2.27, 1.78, 0.42, step, 13, INK, bold=True, font_name="Aptos Display")
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx < len(steps) - 1:
            add_connector(slide, x + 2.08, 2.54, x + 2.38, 2.54)
        x += 2.38
    add_panel(slide, 0.9, 3.78, 5.8, 2.55)
    add_bullets(slide, 1.15, 4.18, 4.95, 1.62, [
        tr(lang, "The system can react when paper stock falls below target.", "Sistem paper stoğu hedefin altına düştüğünde tepki verebilir."),
        tr(lang, "Feedback refresh is part of the refill cycle by default.", "Feedback yenilemesi varsayılan olarak refill döngüsünün parçasıdır."),
        tr(lang, "DergiPark refresh keeps the Turkish pool renewable.", "DergiPark yenilemesi Türkçe havuzu yenilenebilir tutar."),
    ], font_size=13)
    add_panel(slide, 7.0, 3.78, 5.2, 2.55, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.28, 4.08, 4.45, 1.32, tr(lang, "The strongest sentence here is: “The current system is operational, not only experimental: it can refresh feedback, search new papers, and push accepted PDFs back into the annotator pipeline.”", "Buradaki en güçlü cümle şudur: “Mevcut sistem yalnızca deneysel değil, operasyoneldir: geri beslemeyi yenileyebilir, yeni makaleler arayabilir ve kabul edilen PDF'leri tekrar annotator hattına itebilir.”"), 14, INK, bold=True)
    note(
        tr(lang, "Operational loop", "Operasyonel döngü"),
        "Use this slide to prove the project can run as a process, not only as code modules.",
        "Bu slaytı projenin yalnızca kod modülleri değil, çalışan bir süreç olduğunu kanıtlamak için kullanın.",
    )

    # 15. UI interaction
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "AI-relevant UI interaction", "Yapay zeka açısından ilgili arayüz etkileşimi"), tr(lang, "The frontend is not only data entry. It captures structured evidence that the retrieval layer can later learn from.", "Frontend yalnızca veri girişi değildir. Erişim katmanının daha sonra öğrenebileceği yapılandırılmış kanıtı da yakalar."))
    add_small_card(slide, 0.82, 1.62, 3.72, 4.9, tr(lang, "PDF highlight and click resolution", "PDF highlight ve tıklama çözümleme"), tr(lang, "Nutrient names are matched with boundary-aware regex patterns, wrapped in marks, and opened through direct-target, elementsFromPoint, or caret-based fallback when PDF.js text layers behave irregularly.", "Nutrient adları sınır duyarlı regex kalıplarıyla eşleştirilir, mark etiketlerine sarılır ve PDF.js text layer düzensiz davrandığında direct-target, elementsFromPoint veya caret tabanlı fallback ile açılır."), fill_color=GREEN_LIGHT)
    add_small_card(slide, 4.8, 1.62, 3.72, 4.9, tr(lang, "Quick-add popover and structured save", "Hızlı ekleme popover'ı ve yapılandırılmış kayıt"), tr(lang, "A clicked nutrient opens a popover near the highlight, captures value and unit, and adds a structured nutrient row into the current food-item form.", "Tıklanan nutrient, highlight yanında bir popover açar; değer ve birim alır ve mevcut food-item formuna yapılandırılmış nutrient satırı ekler."), fill_color=BLUE_LIGHT)
    add_small_card(slide, 8.78, 1.62, 3.52, 4.9, tr(lang, "Autocomplete ranking and telemetry", "Autocomplete sıralaması ve telemetri"), tr(lang, "Food and nutrient searches use custom ranking rules, and search_sessions record typed queries, shown options, selected options, and abandon cases.", "Food ve nutrient aramaları özel sıralama kuralları kullanır; search_sessions ise yazılan sorguları, gösterilen seçenekleri, seçilen seçenekleri ve terk edilen durumları kaydeder."), fill_color=AMBER_LIGHT)
    note(
        tr(lang, "Why the UI matters to AI", "Arayüz neden yapay zeka için önemlidir"),
        "The UI is where human evidence becomes structured data. The three strongest implementation details are regex-based PDF matching, click fallback logic, and search-session telemetry.",
        "Arayüz, insan kanıtının yapılandırılmış veriye dönüştüğü yerdir. En güçlü üç uygulama detayı regex tabanlı PDF eşleştirme, tıklama fallback mantığı ve search-session telemetrisidir.",
    )

    # 16. Scope boundary
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Current state and scope boundary", "Mevcut durum ve kapsam sınırı"), tr(lang, "This slide is for staying technically strong without overclaiming.", "Bu slayt aşırı iddiaya kaçmadan teknik olarak güçlü kalmak içindir."))
    add_panel(slide, 0.82, 1.62, 5.8, 5.0, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.1, 1.92, 3.0, 0.3, tr(lang, "Implemented now", "Şu anda uygulananlar"), 16, GREEN, bold=True, font_name="Aptos Display")
    add_bullets(slide, 1.08, 2.35, 4.95, 3.7, [
        tr(lang, "Multi-source search across Europe PMC, OpenAlex, Semantic Scholar, and DergiPark", "Europe PMC, OpenAlex, Semantic Scholar ve DergiPark üzerinde çok kaynaklı arama"),
        tr(lang, "Search gate plus stronger metadata filter before acquisition", "Edinimden önce arama kapısı ve daha güçlü meta veri filtresi"),
        tr(lang, "Dual-model embedding support inside metadata scoring", "Meta veri puanlaması içinde çift model embedding desteği"),
        tr(lang, "Batch-updated feedback loop that changes later runs", "Sonraki çalıştırmaları değiştiren toplu güncellenmiş geri besleme döngüsü"),
        tr(lang, "Operational refill loop tied to paper stock", "Paper stoğuna bağlı operasyonel refill döngüsü"),
    ], font_size=14)
    add_panel(slide, 6.85, 1.62, 5.45, 5.0, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.12, 1.92, 3.6, 0.3, tr(lang, "Not claimed at midterm", "Arasınavda iddia edilmeyenler"), 16, AMBER, bold=True, font_name="Aptos Display")
    add_bullets(slide, 7.1, 2.35, 4.65, 3.7, [
        tr(lang, "No trained classifier for query ranking yet", "Henüz query ranking için eğitilmiş sınıflandırıcı yok"),
        tr(lang, "No LLM-based search-term ranking", "LLM tabanlı search-term ranking yok"),
        tr(lang, "No fully automatic extraction pipeline in the current loop", "Mevcut döngüde tam otomatik extraction pipeline yok"),
        tr(lang, "Document segmentation and extraction are deferred to the second semester", "Document segmentation ve extraction ikinci döneme ertelenmiştir"),
    ], font_size=14)
    note(
        tr(lang, "Scope discipline", "Kapsam disiplini"),
        "This is your protection slide. The safest line is: the current system learns through feedback-driven statistical adaptation, but it is not yet a trained classifier.",
        "Bu sizin koruma slaydınızdır. En güvenli cümle şudur: mevcut sistem geri besleme güdümlü istatistiksel uyarlama ile öğreniyor, ancak henüz eğitilmiş bir sınıflandırıcı değil.",
    )

    # 17. Takeaways
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Defense takeaway", "Savunma özeti"), tr(lang, "If you had to compress the whole AI story into a few claims, these are the ones worth defending strongly.", "Tüm yapay zeka hikâyesini birkaç iddiaya sıkıştırmak gerekseydi, güçlü biçimde savunulması gerekenler bunlardır."))
    add_small_card(slide, 0.82, 1.75, 3.8, 3.95, tr(lang, "Claim 1", "İddia 1"), tr(lang, "The system performs staged filtering before acquisition. This is the main reason expert time is not wasted on every retrieved paper.", "Sistem edinimden önce aşamalı filtreleme yapar. Uzman zamanının her getirilen paper üzerinde boşa gitmemesinin temel nedeni budur."), fill_color=GREEN_LIGHT)
    add_small_card(slide, 4.78, 1.75, 3.8, 3.95, tr(lang, "Claim 2", "İddia 2"), tr(lang, "English and Turkish are separate retrieval pools. This is a real technical decision, not only a presentation detail, and DergiPark is part of that strategy.", "İngilizce ve Türkçe ayrı erişim havuzlarıdır. Bu yalnızca sunum detayı değil, gerçek bir teknik karardır ve DergiPark bu stratejinin parçasıdır."), fill_color=BLUE_LIGHT)
    add_small_card(slide, 8.74, 1.75, 3.5, 3.95, tr(lang, "Claim 3", "İddia 3"), tr(lang, "The feedback loop is implemented now. It updates phrases, weighted terms, and source/query priorities for later runs.", "Geri besleme döngüsü şu anda uygulanmıştır. Sonraki çalıştırmalar için phrase'leri, weighted term'leri ve source/query önceliklerini günceller."), fill_color=AMBER_LIGHT)
    add_panel(slide, 1.3, 6.02, 10.75, 0.52)
    add_textbox(slide, 1.55, 6.17, 10.2, 0.2, tr(lang, "If challenged, come back to these three claims and then answer with one concrete implementation detail.", "Zorlanırsanız bu üç iddiaya geri dönün ve sonra tek bir somut uygulama detayıyla cevap verin."), 11, INK, bold=True)
    note(
        tr(lang, "Main defense claims", "Ana savunma iddiaları"),
        "If discussion becomes broad or chaotic, come back to these three claims. Then answer with one concrete implementation detail.",
        "Tartışma dağılırsa bu üç iddiaya geri dönün. Sonra tek bir somut uygulama detayıyla cevap verin.",
    )

    # 18. Appendix divider
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix: implementation detail and Q&A", "Ek: uygulama ayrıntıları ve soru-cevap"), tr(lang, "The next slides are designed for long committee questions.", "Sonraki slaytlar uzun kurul soruları için tasarlanmıştır."), appendix=True)
    add_panel(slide, 0.85, 1.6, 11.55, 4.9, fill_color=AMBER_LIGHT)
    add_textbox(slide, 1.25, 2.05, 10.5, 1.0, tr(lang, "Appendix: implementation detail and Q&A", "Ek: uygulama ayrıntıları ve soru-cevap"), 30, INK, bold=True, font_name="Aptos Display")
    add_textbox(slide, 1.25, 3.18, 9.6, 1.4, tr(lang, "The appendix goes deeper into schema slices, exact thresholds, feedback rules, UI internals, and safe answers for common committee questions.", "Ek bölüm; şema kesitlerine, tam eşiklere, feedback kurallarına, arayüz iç yapısına ve kuruldan gelebilecek yaygın sorular için güvenli cevaplara daha derin iner."), 18, MUTED)
    note(
        tr(lang, "Appendix divider", "Ek ayırıcı"),
        "Use this slide only as a transition. You do not need to speak on it for long.",
        "Bu slaytı yalnızca geçiş için kullanın. Uzun anlatmanıza gerek yok.",
    )

    # 19. Sources and language detail
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix A. Sources and language strategy", "Ek A. Kaynaklar ve dil stratejisi"), tr(lang, "This is the detailed justification for using multiple sources and for not merging EN and TR into one ranking pool.", "Bu, birden fazla kaynak kullanmanın ve EN ile TR'yi tek sıralama havuzunda birleştirmemenin ayrıntılı gerekçesidir."), appendix=True)
    add_panel(slide, 0.82, 1.65, 7.15, 4.95)
    add_bullets(slide, 1.08, 2.0, 6.2, 4.0, [
        tr(lang, "Europe PMC provides open-access biomedical coverage and is the strongest English default surface.", "Europe PMC açık erişimli biyomedikal kapsama sağlar ve İngilizce için en güçlü varsayılan yüzeydir."),
        tr(lang, "OpenAlex broadens metadata coverage and supports both languages when PMC coverage is weak.", "OpenAlex meta veri kapsamını genişletir ve PMC kapsamı zayıfken her iki dili destekler."),
        tr(lang, "Semantic Scholar adds another discovery path and duplicate-aware evidence.", "Semantic Scholar başka bir keşif yolu ve duplicate farkındalıklı kanıt ekler."),
        tr(lang, "DergiPark is the essential Turkish source and justifies the language-specific Turkish workflow.", "DergiPark temel Türkçe kaynaktır ve dile özgü Türkçe iş akışını gerekçelendirir."),
    ], font_size=13)
    add_panel(slide, 8.25, 1.65, 4.05, 4.95, fill_color=GREEN_LIGHT)
    add_bullets(slide, 8.52, 2.0, 3.2, 3.95, [
        tr(lang, "Different source visibility", "Farklı kaynak görünürlüğü"),
        tr(lang, "Different term surfaces and composition phrasing", "Farklı term yüzeyleri ve composition phrasing"),
        tr(lang, "Different source priors and static biases", "Farklı source prior'lar ve statik öncelikler"),
        tr(lang, "Different anchor and query-phrase pools", "Farklı anchor ve query-phrase havuzları"),
    ], font_size=13)
    note(
        tr(lang, "Language strategy detail", "Dil stratejisinin ayrıntısı"),
        "The short answer is that sources, phrasing, and priors differ enough that merging EN and TR would reduce control and likely hurt Turkish coverage.",
        "Kısa cevap şudur: kaynaklar, ifadeler ve öncelikler yeterince farklıdır; EN ve TR'yi birleştirmek kontrolü azaltır ve muhtemelen Türkçe kapsamı zayıflatır.",
    )

    # 20. Exact thresholds
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix B. Exact weights and thresholds", "Ek B. Tam ağırlıklar ve eşikler"), tr(lang, "These are the implemented constants for the two main ranking stages.", "Bunlar iki ana sıralama aşaması için uygulanmış sabitlerdir."), appendix=True)
    add_panel(slide, 0.82, 1.72, 5.75, 4.85, fill_color=GREEN_LIGHT)
    add_textbox(slide, 1.08, 1.98, 3.2, 0.3, tr(lang, "Search gate", "Arama kapısı"), 15, GREEN, bold=True)
    add_textbox(slide, 1.08, 2.4, 4.9, 3.3, tr(lang, "composition phrase        +0.90\nfood term                +0.35\nnutrient term            +0.35\nunit pattern             +0.70\nfood + nutrient combo    +0.45\nmissing metadata         -2.00 reject\nmissing abstract         -1.30 or -0.45\nstrong negatives         up to -2.50\nsoft negatives           up to -1.50\nhealth terms             up to -1.60\npass if score >= -0.35", "composition phrase        +0.90\nfood term                +0.35\nnutrient term            +0.35\nunit pattern             +0.70\nfood + nutrient combo    +0.45\nmissing metadata         -2.00 ret\nmissing abstract         -1.30 veya -0.45\nstrong negatives         en fazla -2.50\nsoft negatives           en fazla -1.50\nhealth terms             en fazla -1.60\ngeçiş için score >= -0.35"), 14, INK, bold=True, font_name="Consolas")
    add_panel(slide, 6.82, 1.72, 5.45, 4.85, fill_color=BLUE_LIGHT)
    add_textbox(slide, 7.08, 1.98, 3.2, 0.3, tr(lang, "Metadata filter", "Meta veri filtresi"), 15, BLUE, bold=True)
    add_textbox(slide, 7.08, 2.4, 4.6, 3.3, tr(lang, "composition phrase        +1.35\nunit pattern             +1.25\nfood hit                 +0.65\nnutrient hit             +0.65\nfood + nutrient combo    +0.75\nembedding above threshold +1.45\nembedding positive bonus +0.75\nsource prior clamp       +/-0.90\nfeedback term clamp      +/-2.50\nfeedback score clamp     +/-6.00\naccept if score >= 1.75", "composition phrase        +1.35\nunit pattern             +1.25\nfood hit                 +0.65\nnutrient hit             +0.65\nfood + nutrient combo    +0.75\nembedding above threshold +1.45\nembedding positive bonus +0.75\nsource prior clamp       +/-0.90\nfeedback term clamp      +/-2.50\nfeedback score clamp     +/-6.00\nkabul için score >= 1.75"), 14, INK, bold=True, font_name="Consolas")
    note(
        tr(lang, "Exact thresholds", "Tam eşikler"),
        "Only use this slide if the committee asks for exact constants. Otherwise summarize the system as additive and threshold-based.",
        "Bu slaytı yalnızca kurul tam sabitleri sorarsa kullanın. Aksi halde sistemi toplamsal ve eşik tabanlı olarak özetleyin.",
    )

    # 21. Feedback scoring detail
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix C. Feedback scoring detail", "Ek C. Feedback puanlama ayrıntısı"), tr(lang, "This slide combines exact label rules with the term-ranking logic used by the feedback updater.", "Bu slayt, feedback updater tarafından kullanılan tam etiket kurallarını ve term ranking mantığını birleştirir."), appendix=True)
    add_panel(slide, 0.82, 1.72, 5.7, 4.85, fill_color=GREEN_LIGHT)
    add_bullets(slide, 1.08, 2.0, 4.85, 2.1, [
        tr(lang, "Positive = latest draft/done with has_data = true and real structured rows", "Pozitif = has_data = true ve gerçek yapılandırılmış satırlara sahip son draft/done"),
        tr(lang, "Negative = global definitely-no-data or at least two unique skips", "Negatif = global definitely-no-data veya en az iki farklı skip"),
        tr(lang, "Conflict = both positive and negative evidence, excluded from learning", "Çatışma = hem pozitif hem negatif kanıt, öğrenmeden çıkarılır"),
    ], font_size=13)
    add_textbox(slide, 1.08, 4.72, 4.95, 1.2, tr(lang, "query_rank = 1.75*title_good + 0.75*title_net + 0.35*ta_net - max(0, title_bad)\nanchor_rank = 1.25*title_good + 1.0*ta_good + 0.5*title_net + 0.25*ta_net - max(0, title_bad)", "query_rank = 1.75*title_good + 0.75*title_net + 0.35*ta_net - max(0, title_bad)\nanchor_rank = 1.25*title_good + 1.0*ta_good + 0.5*title_net + 0.25*ta_net - max(0, title_bad)"), 12, INK, bold=True, font_name="Consolas")
    add_panel(slide, 6.78, 1.72, 5.5, 4.85, fill_color=AMBER_LIGHT)
    add_bullets(slide, 7.04, 2.0, 4.7, 3.9, [
        tr(lang, "n-grams are extracted separately from titles and title+abstract text", "n-gram'ler başlıklardan ve başlık+özetten ayrı ayrı çıkarılır"),
        tr(lang, "positive, negative, and background document frequencies are compared with log-odds style scoring", "pozitif, negatif ve arka plan belge frekansları log-odds benzeri puanlamayla karşılaştırılır"),
        tr(lang, "query_phrases default max = 64, anchor_phrases default max = 16", "query_phrases varsayılan üst sınır = 64, anchor_phrases varsayılan üst sınır = 16"),
        tr(lang, "weighted_terms carries the full scored-term payload into later metadata scoring", "weighted_terms tam scored-term yükünü sonraki meta veri puanlamasına taşır"),
    ], font_size=12)
    note(
        tr(lang, "Search-term ranking answer", "Search-term ranking cevabı"),
        "If asked whether the search terms are ranked by ML or LLM, say no: they are ranked by feedback-driven statistical scoring and hand-designed ranking formulas.",
        "Search term'lerin ML veya LLM ile mi sıralandığını sorarlarsa hayır deyin: feedback güdümlü istatistiksel puanlama ve elle tasarlanmış sıralama formülleriyle sıralanırlar.",
    )

    # 22. Pair/batch/concept feedback
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix D. Pair, batch, and concept feedback", "Ek D. Pair, batch ve concept feedback"), tr(lang, "The feedback exporter learns more than phrases: it also learns which source-query combinations and concept terms produce better outcomes.", "Feedback exporter phrase'lerden fazlasını öğrenir: hangi source-query kombinasyonlarının ve concept term'lerinin daha iyi sonuç ürettiğini de öğrenir."), appendix=True)
    add_small_card(slide, 0.82, 1.78, 3.78, 4.65, tr(lang, "Pair feedback", "Pair feedback"), tr(lang, "Group search hits by source + template_id + source_term.\nScore = novel_positive_count / retrieved.\nAlso derive source_priors as log((positive + 1) / (negative + 1)).", "Search hit'leri source + template_id + source_term ile grupla.\nScore = novel_positive_count / retrieved.\nAyrıca source_priors = log((positive + 1) / (negative + 1)) olarak türetilir."), fill_color=GREEN_LIGHT)
    add_small_card(slide, 4.78, 1.78, 3.78, 4.65, tr(lang, "Batch feedback", "Batch feedback"), tr(lang, "Group by batch_key over paper_search_batches and paper_search_batch_hits.\nScore = positive_count / retrieved.", "paper_search_batches ve paper_search_batch_hits üzerinde batch_key ile grupla.\nScore = positive_count / retrieved."), fill_color=BLUE_LIGHT)
    add_small_card(slide, 8.74, 1.78, 3.55, 4.65, tr(lang, "Concept feedback", "Concept feedback"), tr(lang, "Group by source_term.\nScore = (positive_count - negative_count) / retrieved.\nThis later reorders the concept pool.", "source_term ile grupla.\nScore = (positive_count - negative_count) / retrieved.\nBu daha sonra concept havuzunu yeniden sıralar."), fill_color=AMBER_LIGHT)
    note(
        tr(lang, "Non-phrase feedback", "Phrase dışı feedback"),
        "If asked what learns besides terms, answer: pair scores, batch scores, concept scores, and source priors.",
        "Terimler dışında neyin öğrendiğini sorarlarsa şöyle cevap verin: pair score'lar, batch score'lar, concept score'lar ve source prior'lar.",
    )

    # 23. Acquisition and validation
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix E. Acquisition and validation", "Ek E. Edinim ve doğrulama"), tr(lang, "Even after metadata acceptance, papers still need a robust PDF path and a full-text validation step before becoming accepted records.", "Meta veri kabulünden sonra bile, paper'ların kabul edilmiş kayda dönüşmeden önce sağlam bir PDF yoluna ve tam metin doğrulama adımına ihtiyacı vardır."), appendix=True)
    add_panel(slide, 0.82, 1.68, 11.45, 1.28, fill_color=GREEN_LIGHT)
    steps = [tr(lang, "dedupe", "dedupe"), tr(lang, "metadata accept", "metadata kabul"), tr(lang, "PDF fetch", "PDF fetch"), tr(lang, "pdftotext validation", "pdftotext doğrulama"), tr(lang, "manifest + upload", "manifest + upload")]
    x = 1.02
    for idx, step in enumerate(steps):
        add_panel(slide, x, 2.0, 1.8, 0.62, fill_color=PANEL if idx % 2 else BLUE_LIGHT)
        frame = add_textbox(slide, x + 0.1, 2.17, 1.6, 0.18, step, 11, INK, bold=True)
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx < len(steps) - 1:
            add_connector(slide, x + 1.8, 2.31, x + 2.08, 2.31)
        x += 2.08
    add_panel(slide, 0.82, 3.42, 5.6, 3.1)
    add_bullets(slide, 1.08, 3.82, 4.8, 2.0, [
        tr(lang, "Try OA XML and direct PDF links first.", "Önce OA XML ve doğrudan PDF bağlantılarını dene."),
        tr(lang, "Use package extraction or HTML-discovered links when direct links fail.", "Doğrudan bağlantı başarısız olursa paket çıkarımı veya HTML'den keşfedilen bağlantıları kullan."),
        tr(lang, "Use curl / fallback paths for unstable source behavior.", "Kararsız kaynak davranışı için curl / fallback yollarını kullan."),
    ], font_size=13)
    add_panel(slide, 6.72, 3.42, 5.55, 3.1, fill_color=AMBER_LIGHT)
    add_bullets(slide, 6.98, 3.82, 4.75, 2.0, [
        tr(lang, "Extract full text with pdftotext.", "Tam metni pdftotext ile çıkar."),
        tr(lang, "Check for composition framing, table/method evidence, nutrient-unit signals, and food markers.", "Composition çerçevesini, tablo/metot kanıtını, nutrient-birim sinyallerini ve food işaretlerini kontrol et."),
        tr(lang, "Reject candidates that still fail full-text usefulness checks.", "Tam metin yararlılık kontrollerini yine de geçemeyen adayları reddet."),
    ], font_size=13)
    note(
        tr(lang, "Acquisition path", "Edinim yolu"),
        "If asked how you know a paper is really usable, answer: metadata accept is not enough; the project fetches the PDF and validates the full text before final acceptance.",
        "Bir paper'ın gerçekten kullanılabilir olduğunu nasıl bildiğinizi sorarlarsa şöyle cevap verin: meta veri kabulü yeterli değildir; proje PDF'yi getirir ve son kabulden önce tam metni doğrular.",
    )

    # 24. UI internals
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix F. UI internals relevant to AI", "Ek F. Yapay zeka açısından ilgili arayüz iç yapısı"), tr(lang, "These details matter when the committee asks how human interaction is turned into structured signals.", "Kurul insan etkileşiminin nasıl yapılandırılmış sinyallere dönüştüğünü sorduğunda bu ayrıntılar önemlidir."), appendix=True)
    add_small_card(slide, 0.82, 1.72, 3.7, 4.85, tr(lang, "PDF highlight", "PDF highlight"), tr(lang, "Build boundary-aware nutrient regex patterns, scan PDF.js text spans after render, wrap matches in mark tags, and resolve overlaps so shorter nutrient names do not break longer ones.", "Sınır duyarlı nutrient regex kalıpları oluştur, render sonrası PDF.js text span'lerini tara, eşleşmeleri mark etiketlerine sar ve kısa nutrient adlarının uzun olanları bozmaması için çakışmaları çöz."), fill_color=GREEN_LIGHT)
    add_small_card(slide, 4.8, 1.72, 3.72, 4.85, tr(lang, "Click resolution", "Tıklama çözümleme"), tr(lang, "Try the direct event target first, then fall back to elementsFromPoint, then to caretPositionFromPoint or caretRangeFromPoint because PDF.js text layers do not behave like normal HTML.", "Önce doğrudan event target'ı dene, sonra elementsFromPoint'e, sonra caretPositionFromPoint veya caretRangeFromPoint'a dön; çünkü PDF.js text layer'ları normal HTML gibi davranmaz."), fill_color=BLUE_LIGHT)
    add_small_card(slide, 8.8, 1.72, 3.48, 4.85, tr(lang, "Autocomplete and telemetry", "Autocomplete ve telemetri"), tr(lang, "Food and nutrient autocomplete use custom ranking rules, and search_sessions logs typed queries, shown options, selected options, status, and timestamps.", "Food ve nutrient autocomplete özel sıralama kuralları kullanır; search_sessions ise yazılan sorguları, gösterilen seçenekleri, seçilen seçenekleri, durumu ve zaman damgalarını kaydeder."), fill_color=AMBER_LIGHT)
    note(
        tr(lang, "UI detail", "Arayüz ayrıntısı"),
        "The best AI-facing sentence here is: the UI is where human evidence becomes structured data and telemetry.",
        "Buradaki en iyi yapay zeka odaklı cümle şudur: arayüz, insan kanıtının yapılandırılmış veriye ve telemetriye dönüştüğü yerdir.",
    )

    # 25. Run audit
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix G. Audited sample run summary", "Ek G. Denetlenmiş örnek çalıştırma özeti"), tr(lang, "This slide intentionally separates stage counts from recorded acquisition outcomes so the sample run is not misread as a strict funnel.", "Bu slayt, örnek çalıştırmanın katı bir huni gibi yanlış okunmaması için aşama sayılarını kayıtlı edinim çıktılarından bilinçli olarak ayırır."), appendix=True)
    add_panel(slide, 0.82, 1.72, 5.82, 4.85)
    add_textbox(slide, 1.08, 1.98, 3.6, 0.3, tr(lang, "Stage counts (TR live run 2026-03-30)", "Aşama sayıları (TR canlı çalıştırma 2026-03-30)"), 14, GREEN, bold=True)
    add_hbar(slide, tr(lang, "Raw hits", "Ham hit"), run_summary["hits"], run_summary["hits"], 1.08, 2.5, 2.65, BLUE)
    add_hbar(slide, tr(lang, "Passed search gate", "Arama kapısını geçti"), run_summary["search_gate_pass"], run_summary["hits"], 1.08, 3.0, 2.65, GREEN)
    add_hbar(slide, tr(lang, "Passed metadata filter", "Meta veri filtresini geçti"), run_summary["metadata_pass"], run_summary["hits"], 1.08, 3.5, 2.65, AMBER)
    add_textbox(slide, 1.08, 4.24, 4.95, 1.5, tr(lang, "These counts come from the manifest summary and are valid as stage counts up to metadata_pass.", "Bu sayılar manifest özetinden gelir ve metadata_pass aşamasına kadar aşama sayıları olarak geçerlidir."), 12, MUTED)
    add_panel(slide, 6.9, 1.72, 5.38, 4.85, fill_color=AMBER_LIGHT)
    add_textbox(slide, 7.16, 1.98, 3.6, 0.3, tr(lang, "Recorded acquisition outcomes", "Kayıtlı edinim çıktıları"), 14, AMBER, bold=True)
    outcome_max = max(run_summary["metadata_pass"], run_summary["terminal_recorded"])
    add_hbar(slide, tr(lang, "PDF fetch failed", "PDF fetch başarısız"), run_summary["pdf_fetch_fail"], outcome_max, 7.15, 2.5, 2.35, AMBER)
    add_hbar(slide, tr(lang, "PDF validation failed", "PDF doğrulama başarısız"), run_summary["pdf_validation_fail"], outcome_max, 7.15, 3.0, 2.35, RED)
    add_hbar(slide, tr(lang, "Accepted", "Kabul edildi"), run_summary["accepted"], outcome_max, 7.15, 3.5, 2.35, GREEN)
    add_hbar(slide, tr(lang, "No terminal result in manifest", "Manifestte terminal sonuç yok"), run_summary["unaccounted"], outcome_max, 7.15, 4.0, 2.35, BLUE)
    add_textbox(slide, 7.16, 4.72, 4.5, 0.9, tr(lang, "Important audit note: 78 metadata-pass candidates exist, but only 66 have recorded terminal results in this manifest. Therefore this sample must not be presented as a strict funnel.", "Önemli denetim notu: Bu manifestte 78 metadata-pass aday vardır, ancak bunların yalnızca 66'sının kayıtlı terminal sonucu vardır. Bu nedenle bu örnek katı bir huni olarak sunulmamalıdır."), 11, INK, bold=True)
    note(
        tr(lang, "Run-summary caveat", "Çalıştırma özeti uyarısı"),
        "Use this slide to show rigor, not to defend raw counts. Stage counts are real up to metadata_pass, but this manifest leaves 12 metadata-pass candidates without terminal outcomes.",
        "Bu slaytı ham sayıları savunmak için değil, titizlik göstermek için kullanın. Aşama sayıları metadata_pass'e kadar gerçektir, ancak bu manifest 12 metadata-pass adayı terminal sonuç olmadan bırakır.",
    )

    # 26. Q&A
    slide_no += 1
    slide = new_slide(prs, lang, slide_no, tr(lang, "Appendix H. Likely questions and safe answers", "Ek H. Olası sorular ve güvenli cevaplar"), tr(lang, "These are the short answers to return to under pressure.", "Baskı altında geri dönülmesi gereken kısa cevaplar bunlardır."), appendix=True)
    add_panel(slide, 0.82, 1.72, 5.78, 4.85)
    add_textbox(slide, 1.08, 1.98, 5.0, 3.8, tr(lang, "Q: Is the feedback loop implemented?\nA: Yes. It is batch-updated, not online, and it already changes later crawler runs.\n\nQ: Do you use an LLM to rank search terms?\nA: No. Search-term ranking is feedback-driven statistical scoring, not LLM ranking.\n\nQ: Did you train a model for query ranking?\nA: Not yet. The current system learns through feedback statistics and embeddings, not a trained classifier.", "S: Geri besleme döngüsü uygulanmış mı?\nC: Evet. Toplu güncelleniyor, çevrim içi değil ve sonraki crawler çalıştırmalarını şimdiden değiştiriyor.\n\nS: Search term'leri sıralamak için LLM kullanıyor musunuz?\nC: Hayır. Search-term ranking, LLM değil, feedback güdümlü istatistiksel puanlamadır.\n\nS: Query ranking için model eğittiniz mi?\nC: Henüz değil. Mevcut sistem, eğitilmiş sınıflandırıcıdan çok feedback istatistikleri ve embedding'lerle öğreniyor."), 14, INK)
    add_panel(slide, 6.9, 1.72, 5.38, 4.85, fill_color=GREEN_LIGHT)
    add_textbox(slide, 7.16, 1.98, 4.65, 3.8, tr(lang, "Q: Why not download everything?\nA: Because metadata filtering is much cheaper and protects expert time.\n\nQ: Why separate English and Turkish?\nA: Because source coverage, terminology, and source priors differ.\n\nQ: Why not only embeddings?\nA: Domain cues like units, composition phrases, and food/nutrient combinations are too valuable to ignore.", "S: Neden her şeyi indirmiyorsunuz?\nC: Çünkü meta veri filtreleme çok daha ucuzdur ve uzman zamanını korur.\n\nS: Neden İngilizce ve Türkçeyi ayırıyorsunuz?\nC: Çünkü kaynak kapsamı, terminoloji ve kaynak öncelikleri farklıdır.\n\nS: Neden yalnızca embedding değil?\nC: Birimler, composition phrase'ler ve food/nutrient kombinasyonları gibi alan ipuçları görmezden gelinemeyecek kadar değerlidir."), 14, INK)
    note(
        tr(lang, "Final Q&A slide", "Son Soru-Cevap slaytı"),
        "The three safest sentences are: the feedback loop is implemented but batch-updated; search-term ranking is statistical not LLM-based; and the ranking layer is additive.", 
        "En güvenli üç cümle şunlardır: geri besleme döngüsü uygulanmıştır ama toplu güncellenir; search-term ranking istatistikseldir ve LLM tabanlı değildir; sıralama katmanı toplamsaldır.",
    )

    pptx_name = "OpenNutri_AI_Algorithm_Master_Deck_EN.pptx" if lang == "en" else "OpenNutri_AI_Algorithm_Master_Deck_TR.pptx"
    notes_name = "OpenNutri_AI_Algorithm_Master_Deck_EN.md" if lang == "en" else "OpenNutri_AI_Algorithm_Master_Deck_TR.md"
    pptx_path = PPTX_DIR / pptx_name
    notes_path = NOTES_DIR / notes_name
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    NOTES_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    write_notes(notes_path, deck_title, notes, lang=lang)
    return pptx_path, notes_path


def main() -> None:
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    PDF_DIR.mkdir(parents=True, exist_ok=True)
    NOTES_DIR.mkdir(parents=True, exist_ok=True)
    for lang in ("en", "tr"):
        pptx_path, notes_path = build_master_deck(lang)
        export_pdf(pptx_path, PDF_DIR)
        print(f"Created {pptx_path}")
        print(f"Created {notes_path}")


if __name__ == "__main__":
    main()
